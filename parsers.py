import os
import io
from pathlib import Path
from typing import Dict, List, Optional, Any
import logging
import requests
import tempfile

logger = logging.getLogger(__name__)

# Default minimum image size (width, height) - skip smaller images
MIN_IMAGE_SIZE = (50, 50)


class DocumentParser:
    """Universal document parser supporting multiple parsing engines."""
    
    SUPPORTED_EXTENSIONS = {
        'pdf': ['pdf'],
        'word': ['docx', 'doc'],
        'powerpoint': ['pptx', 'ppt'],
        'excel': ['xlsx', 'xls', 'csv'],
        'text': ['txt', 'md', 'html', 'htm']
    }
    
    def __init__(
        self,
        tool: Optional[str] = None,
        detect_scanned: bool = True,
        extract_images: bool = False,
        ocr_lang: str = 'eng'
    ):
        """
        Initialize document parser.
        
        Args:
            tool: Specific tool to use (None for auto-detection)
            detect_scanned: Whether to detect scanned/hybrid PDFs
            extract_images: Whether to extract images
            ocr_lang: OCR language code
        """
        self.tool = tool
        self.detect_scanned = detect_scanned
        self.extract_images = extract_images
        self.ocr_lang = ocr_lang
        
    def parse(self, document_path: str) -> Dict[str, Any]:
        """
        Parse document and return structured content.
        
        Args:
            document_path: Path to document or URL
            
        Returns:
            Dictionary containing parsed content and metadata
        """
        # Handle URL downloads
        if document_path.startswith('http://') or document_path.startswith('https://'):
            document_path = self._download_from_url(document_path)
        
        # Verify file exists
        if not os.path.exists(document_path):
            raise FileNotFoundError(f"Document not found: {document_path}")
        
        # Detect file type
        file_ext = Path(document_path).suffix.lower().lstrip('.')
        doc_type = self._get_document_type(file_ext)
        
        # Select appropriate parser
        if self.tool:
            parser_method = self._get_parser_method(self.tool)
        else:
            parser_method = self._auto_select_parser(doc_type, document_path)
        
        # Parse document
        logger.info(f"Parsing {document_path} using {parser_method.__name__}")
        results = parser_method(document_path)
        
        # Add metadata
        results['file_name'] = Path(document_path).name
        results['file_type'] = doc_type
        results['file_size'] = os.path.getsize(document_path)
        
        return results
    
    def _download_from_url(self, url: str) -> str:
        """Download document from URL to temporary file."""
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            # Determine file extension from URL or content-type
            ext = Path(url).suffix or '.pdf'
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
                tmp_file.write(response.content)
                return tmp_file.name
        except Exception as e:
            raise Exception(f"Failed to download from URL: {str(e)}")
    
    def _get_document_type(self, extension: str) -> str:
        """Determine document type from extension."""
        for doc_type, exts in self.SUPPORTED_EXTENSIONS.items():
            if extension in exts:
                return doc_type
        return 'unknown'
    
    def _auto_select_parser(self, doc_type: str, path: str) -> callable:
        """Auto-select best parser based on document type."""
        if doc_type == 'pdf':
            # Check if scanned PDF
            if self.detect_scanned:
                is_scanned = self._is_scanned_pdf(path)
                if is_scanned:
                    logger.info("Detected scanned PDF, using OCR")
                    return self._parse_with_tesseract
            
            # Default to PyMuPDF4LLM for regular PDFs
            return self._parse_with_pymupdf4llm
        
        elif doc_type == 'word':
            return self._parse_word
        elif doc_type == 'powerpoint':
            return self._parse_powerpoint
        elif doc_type == 'excel':
            return self._parse_excel
        elif doc_type == 'text':
            return self._parse_text
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
    
    def _get_parser_method(self, tool_name: str) -> callable:
        """Get parser method by tool name."""
        method_map = {
            'pymupdf': self._parse_with_pymupdf,
            'pymupdf4llm': self._parse_with_pymupdf4llm,
            'docling': self._parse_with_docling,
            'unstructured': self._parse_with_unstructured,
            'pdfplumber': self._parse_with_pdfplumber,
            'python-docx': self._parse_word,
            'python-pptx': self._parse_powerpoint,
            'openpyxl': self._parse_excel,
            'tesseract_ocr': self._parse_with_tesseract,
            'easyocr': self._parse_with_easyocr
        }
        
        if tool_name not in method_map:
            raise ValueError(f"Unknown tool: {tool_name}")
        
        return method_map[tool_name]
    
    def _is_scanned_pdf(self, pdf_path: str) -> bool:
        """Detect if PDF is scanned."""
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            pages_without_text = 0
            
            # Sample first 5 pages
            for page_num in range(min(5, total_pages)):
                page = doc[page_num]
                text = page.get_text().strip()
                if not text or len(text) < 50:
                    pages_without_text += 1
            
            doc.close()
            
            # If most sampled pages have no text, it's likely scanned
            return pages_without_text / min(5, total_pages) > 0.7
        except Exception as e:
            logger.warning(f"Error detecting scanned PDF: {e}")
            return False
    
    def _analyze_pdf(self, pdf_path: str) -> Dict[str, Any]:
        """Analyze PDF structure and type."""
        try:
            import fitz
            
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            
            scanned_pages = []
            hybrid_pages = []
            digital_pages = []
            
            for page_num in range(total_pages):
                page = doc[page_num]
                text = page.get_text().strip()
                images = page.get_images()
                
                has_text = len(text) > 50
                has_images = len(images) > 0
                
                if not has_text and has_images:
                    scanned_pages.append(page_num + 1)
                elif has_text and has_images:
                    hybrid_pages.append(page_num + 1)
                elif has_text:
                    digital_pages.append(page_num + 1)
            
            doc.close()
            
            # Determine overall type
            if len(scanned_pages) == total_pages:
                pdf_type = "Scanned"
            elif len(digital_pages) == total_pages:
                pdf_type = "Digital"
            else:
                pdf_type = "Hybrid"
            
            return {
                'type': pdf_type,
                'total_pages': total_pages,
                'scanned_pages': scanned_pages,
                'hybrid_pages': hybrid_pages,
                'digital_pages': digital_pages,
                'has_text': len(digital_pages) > 0 or len(hybrid_pages) > 0
            }
        except Exception as e:
            logger.warning(f"Error analyzing PDF: {e}")
            return {}
    
    def _parse_with_pymupdf(self, pdf_path: str) -> Dict[str, Any]:
        """Parse PDF using PyMuPDF."""
        import fitz
        
        doc = fitz.open(pdf_path)
        content = []
        pages = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            pages.append({
                'page_number': page_num + 1,
                'content': text,
                'metadata': {
                    'width': page.rect.width,
                    'height': page.rect.height
                }
            })
            
            content.append(f"## Page {page_num + 1}\n\n{text}\n")
        
        images = []
        if self.extract_images:
            images = self._extract_images_pymupdf(doc)
        
        doc.close()
        
        pdf_analysis = None
        if self.detect_scanned:
            pdf_analysis = self._analyze_pdf(pdf_path)
        
        return {
            'tool_used': 'PyMuPDF',
            'content': '\n'.join(content),
            'pages': pages,
            'images': images,
            'metadata': {'page_count': len(pages)},
            'pdf_analysis': pdf_analysis
        }
    
    def _parse_with_pymupdf4llm(self, pdf_path: str) -> Dict[str, Any]:
        """Parse PDF using PyMuPDF4LLM for LLM-optimized output."""
        try:
            import pymupdf4llm
            
            md_text = pymupdf4llm.to_markdown(pdf_path)
            
            # Also get page-by-page for visualization
            import fitz
            doc = fitz.open(pdf_path)
            pages = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                
                # Get bounding boxes
                bboxes = []
                for block in page.get_text("dict")["blocks"]:
                    if block.get("type") == 0:  # Text block
                        bboxes.append({
                            'type': 'text',
                            'bbox': block["bbox"]
                        })
                    elif block.get("type") == 1:  # Image block
                        bboxes.append({
                            'type': 'image',
                            'bbox': block["bbox"]
                        })
                
                pages.append({
                    'page_number': page_num + 1,
                    'content': text,
                    'bboxes': bboxes,
                    'metadata': {
                        'width': page.rect.width,
                        'height': page.rect.height
                    }
                })
            
            images = []
            if self.extract_images:
                images = self._extract_images_pymupdf(doc)
            
            doc.close()
            
            pdf_analysis = None
            if self.detect_scanned:
                pdf_analysis = self._analyze_pdf(pdf_path)
            
            return {
                'tool_used': 'PyMuPDF4LLM',
                'content': md_text,
                'pages': pages,
                'images': images,
                'metadata': {'page_count': len(pages)},
                'pdf_analysis': pdf_analysis
            }
        except ImportError:
            logger.warning("pymupdf4llm not available, falling back to PyMuPDF")
            return self._parse_with_pymupdf(pdf_path)
    
    def _parse_with_pdfplumber(self, pdf_path: str) -> Dict[str, Any]:
        """Parse PDF using pdfplumber (better for tables)."""
        try:
            import pdfplumber
            
            content = []
            pages = []
            
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    
                    # Extract tables
                    tables = page.extract_tables()
                    table_md = ""
                    
                    for table_idx, table in enumerate(tables):
                        if table:
                            table_md += f"\n**Table {table_idx + 1}:**\n\n"
                            # Convert to markdown table
                            if table[0]:
                                table_md += "| " + " | ".join(str(cell) for cell in table[0]) + " |\n"
                                table_md += "|" + "|".join(["---"] * len(table[0])) + "|\n"
                                for row in table[1:]:
                                    table_md += "| " + " | ".join(str(cell) if cell else "" for cell in row) + " |\n"
                            table_md += "\n"
                    
                    page_content = f"{text}\n{table_md}"
                    
                    pages.append({
                        'page_number': page_num + 1,
                        'content': page_content,
                        'metadata': {
                            'width': page.width,
                            'height': page.height,
                            'tables_count': len(tables)
                        }
                    })
                    
                    content.append(f"## Page {page_num + 1}\n\n{page_content}\n")
            
            pdf_analysis = None
            if self.detect_scanned:
                pdf_analysis = self._analyze_pdf(pdf_path)
            
            return {
                'tool_used': 'pdfplumber',
                'content': '\n'.join(content),
                'pages': pages,
                'images': [],
                'metadata': {'page_count': len(pages)},
                'pdf_analysis': pdf_analysis
            }
        except ImportError:
            logger.warning("pdfplumber not available, falling back to PyMuPDF")
            return self._parse_with_pymupdf(pdf_path)
    
    def _parse_with_tesseract(self, pdf_path: str) -> Dict[str, Any]:
        """Parse scanned PDF using Tesseract OCR."""
        try:
            import pytesseract
            from pdf2image import convert_from_path
            
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=300)
            
            content = []
            pages = []
            
            for page_num, image in enumerate(images):
                # Perform OCR
                text = pytesseract.image_to_string(image, lang=self.ocr_lang)
                
                pages.append({
                    'page_number': page_num + 1,
                    'content': text,
                    'metadata': {
                        'ocr_confidence': 'N/A',
                        'width': image.width,
                        'height': image.height
                    }
                })
                
                content.append(f"## Page {page_num + 1}\n\n{text}\n")
            
            return {
                'tool_used': 'Tesseract OCR',
                'content': '\n'.join(content),
                'pages': pages,
                'images': [],
                'metadata': {'page_count': len(pages)},
                'pdf_analysis': {'type': 'Scanned', 'total_pages': len(pages)}
            }
        except ImportError as e:
            raise ImportError(f"Tesseract OCR dependencies not available: {e}")
    
    def _parse_with_easyocr(self, pdf_path: str) -> Dict[str, Any]:
        """Parse scanned PDF using EasyOCR."""
        try:
            import easyocr
            from pdf2image import convert_from_path
            import numpy as np
            
            # Initialize EasyOCR reader
            reader = easyocr.Reader([self.ocr_lang])
            
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=300)
            
            content = []
            pages = []
            
            for page_num, image in enumerate(images):
                # Convert PIL to numpy array
                img_array = np.array(image)
                
                # Perform OCR
                results = reader.readtext(img_array)
                text = '\n'.join([result[1] for result in results])
                
                pages.append({
                    'page_number': page_num + 1,
                    'content': text,
                    'metadata': {
                        'detections': len(results),
                        'width': image.width,
                        'height': image.height
                    }
                })
                
                content.append(f"## Page {page_num + 1}\n\n{text}\n")
            
            return {
                'tool_used': 'EasyOCR',
                'content': '\n'.join(content),
                'pages': pages,
                'images': [],
                'metadata': {'page_count': len(pages)},
                'pdf_analysis': {'type': 'Scanned', 'total_pages': len(pages)}
            }
        except ImportError as e:
            raise ImportError(f"EasyOCR dependencies not available: {e}")
        
    def _parse_with_docling(self, pdf_path: str) -> Dict[str, Any]:
        """Parse PDF using Docling and return full + page-by-page markdown."""
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption

            # Important: batch_size=1 avoids tensor padding errors
            pipeline_options = PdfPipelineOptions(
                layout_batch_size=1,
                table_batch_size=1,
                ocr_batch_size=1,
            )

            converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )

            result = converter.convert(pdf_path, raises_on_error=False)

            # If docling reports failure, return errors (so Streamlit can show them)
            if getattr(result, "status", None) and str(result.status).endswith("FAILURE"):
                return {
                    "tool_used": "Docling",
                    "content": "",
                    "pages": [],
                    "images": [],
                    "metadata": {},
                    "pdf_analysis": None,
                    "errors": [str(e) for e in (result.errors or [])],
                }

            doc = result.document

            # --- Full markdown (whole doc)
            full_md = doc.export_to_markdown(page_break_placeholder="\n\n---\n\n")

            # --- Page-by-page markdown
            pages = []
            page_count = doc.num_pages()  # DoclingDocument.num_pages() exists in docs
            for i in range(page_count):
                page_md = doc.export_to_markdown(page_no=i)
                pages.append({
                    "page_number": i + 1,
                    "content": page_md,
                    "metadata": {}
                })

            # If full_md is empty (rare bug), fallback to plain text export
            if not full_md.strip():
                full_md = doc.export_to_text()

            return {
                "tool_used": "Docling",
                "content": full_md,
                "pages": pages,
                "images": [],
                "metadata": {"page_count": page_count},
                "pdf_analysis": None,
            }

        except ImportError:
            logger.warning("Docling not available, falling back to PyMuPDF")
            return self._parse_with_pymupdf(pdf_path)
        except Exception as e:
            logger.exception("Docling failed, falling back to PyMuPDF")
            return self._parse_with_pymupdf(pdf_path)

    
    def _parse_with_unstructured(self, pdf_path: str) -> Dict[str, Any]:
        """Parse document using Unstructured."""
        try:
            from unstructured.partition.auto import partition
            
            elements = partition(filename=pdf_path)
            content = '\n\n'.join([str(el) for el in elements])
            
            return {
                'tool_used': 'Unstructured',
                'content': content,
                'pages': [],
                'images': [],
                'metadata': {'elements_count': len(elements)},
                'pdf_analysis': None
            }
        except ImportError:
            logger.warning("Unstructured not available, falling back to PyMuPDF")
            return self._parse_with_pymupdf(pdf_path)
    
    def _parse_word(self, doc_path: str) -> Dict[str, Any]:
        """Parse Word document."""
        try:
            from docx import Document
            
            doc = Document(doc_path)
            content = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                content.append("\n**Table:**\n")
                for row in table.rows:
                    content.append(" | ".join([cell.text for cell in row.cells]))
            
            full_content = '\n\n'.join(content)
            
            return {
                'tool_used': 'python-docx',
                'content': full_content,
                'pages': [{'page_number': 1, 'content': full_content}],
                'images': [],
                'metadata': {
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables)
                }
            }
        except ImportError:
            raise ImportError("python-docx not available")
    
    def _parse_powerpoint(self, ppt_path: str) -> Dict[str, Any]:
        """Parse PowerPoint presentation."""
        try:
            from pptx import Presentation
            
            prs = Presentation(ppt_path)
            content = []
            pages = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                slide_content = '\n'.join(slide_text)
                
                pages.append({
                    'page_number': slide_num + 1,
                    'content': slide_content
                })
                
                content.append(f"## Slide {slide_num + 1}\n\n{slide_content}\n")
            
            return {
                'tool_used': 'python-pptx',
                'content': '\n'.join(content),
                'pages': pages,
                'images': [],
                'metadata': {'slides': len(prs.slides)}
            }
        except ImportError:
            raise ImportError("python-pptx not available")
    
    def _parse_excel(self, excel_path: str) -> Dict[str, Any]:
        """Parse Excel spreadsheet."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(excel_path, data_only=True)
            content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content.append(f"## Sheet: {sheet_name}\n")
                
                # Convert to markdown table
                for row in sheet.iter_rows(values_only=True):
                    row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    content.append(f"| {row_str} |")
                
                content.append("\n")
            
            return {
                'tool_used': 'openpyxl',
                'content': '\n'.join(content),
                'pages': [],
                'images': [],
                'metadata': {'sheets': len(wb.sheetnames)}
            }
        except ImportError:
            raise ImportError("openpyxl not available")
    
    def _parse_text(self, text_path: str) -> Dict[str, Any]:
        """Parse text file."""
        with open(text_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            'tool_used': 'Built-in',
            'content': content,
            'pages': [{'page_number': 1, 'content': content}],
            'images': [],
            'metadata': {'lines': len(content.split('\n'))}
        }
    
    def _extract_images_pymupdf(self, doc) -> List[Dict[str, Any]]:
        """Extract images from PyMuPDF document."""
        from PIL import Image
        import io
        
        images = []
        extracted_xrefs = set()  # Avoid duplicates
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    
                    # Skip if already extracted (some images appear on multiple pages)
                    if xref in extracted_xrefs:
                        continue
                    
                    extracted_xrefs.add(xref)
                    
                    # Extract image
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Convert to PIL Image
                    try:
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        
                        # Convert to RGB if necessary (fixes some display issues)
                        if pil_image.mode not in ('RGB', 'RGBA'):
                            if pil_image.mode == 'CMYK':
                                pil_image = pil_image.convert('RGB')
                            elif pil_image.mode == '1':  # Binary
                                pil_image = pil_image.convert('L')  # Grayscale
                            elif pil_image.mode == 'P':  # Palette
                                pil_image = pil_image.convert('RGB')
                        
                        # Skip very small images (likely artifacts)
                        if pil_image.width < MIN_IMAGE_SIZE[0] or pil_image.height < MIN_IMAGE_SIZE[1]:
                            continue
                        
                        images.append({
                            'page': page_num + 1,
                            'index': img_index,
                            'image': pil_image,
                            'ext': image_ext,
                            'width': pil_image.width,
                            'height': pil_image.height,
                            'mode': pil_image.mode
                        })
                        
                    except Exception as img_error:
                        logger.warning(f"Failed to convert image to PIL format: {img_error}")
                        # Try alternative method: extract from page rendering
                        try:
                            bbox = img[1:5]  # Get bounding box
                            pix = page.get_pixmap(clip=bbox, matrix=fitz.Matrix(2, 2))
                            img_data = pix.tobytes("png")
                            pil_image = Image.open(io.BytesIO(img_data))
                            
                            if pil_image.width >= MIN_IMAGE_SIZE[0] and pil_image.height >= MIN_IMAGE_SIZE[1]:
                                images.append({
                                    'page': page_num + 1,
                                    'index': img_index,
                                    'image': pil_image,
                                    'ext': 'png',
                                    'width': pil_image.width,
                                    'height': pil_image.height,
                                    'mode': pil_image.mode
                                })
                        except:
                            pass
                        
                except Exception as e:
                    logger.warning(f"Failed to extract image on page {page_num + 1}: {e}")
                    continue
        
        return images